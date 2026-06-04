"""
Think! Ventures Foundation -- Board Recruitment Presentation
Built on the official Think! Ventures .potx template
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from zipfile import ZipFile
import shutil, os

# ─── PATHS ───
TEMPLATE_SRC = r'C:\Users\Chris\Desktop\WEBSITES\THINK! VENTURES\docs\Think_Ventures_Template.potx'
TEMPLATE      = r'C:\Users\Chris\Desktop\WEBSITES\THINK! VENTURES\docs\Think_Ventures_Template_converted.pptx'
OUTPUT        = r'C:\Users\Chris\Downloads\Think_Ventures_Board_FINAL4.pptx'
BELLA         = r'C:\Users\Chris\Desktop\WEBSITES\THINK! VENTURES\assets\images\bella-mascot.png'

# ─── BRAND COLORS ───
TEAL       = RGBColor(0x0D, 0x4F, 0x4F)
GOLD       = RGBColor(0xF5, 0xA6, 0x23)
EMERALD    = RGBColor(0x10, 0xB9, 0x81)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT  = RGBColor(0x1A, 0x2A, 0x3A)
MUTED_TEXT = RGBColor(0x4A, 0x5A, 0x6A)
CARD_BG    = RGBColor(0xE8, 0xED, 0xF2)

# ─── CONVERT .potx TO .pptx ───
if not os.path.exists(TEMPLATE):
    shutil.copy2(TEMPLATE_SRC, TEMPLATE)
    with ZipFile(TEMPLATE, 'r') as zin:
        all_data = {n: zin.read(n) for n in zin.namelist()}
    ct = all_data['[Content_Types].xml'].decode('utf-8').replace(
        'presentationml.template.main+xml',
        'presentationml.presentation.main+xml')
    all_data['[Content_Types].xml'] = ct.encode('utf-8')
    with ZipFile(TEMPLATE, 'w') as zout:
        for n, d in all_data.items():
            zout.writestr(n, d)

prs = Presentation(TEMPLATE)

# Layout references
LO_TITLE   = prs.slide_layouts[0]  # Title Slide
LO_CONTENT = prs.slide_layouts[1]  # Title and Content
LO_SECTION = prs.slide_layouts[2]  # Section Header
LO_BLANK   = prs.slide_layouts[3]  # Blank

# ─── HELPERS ───

def txt(slide, left, top, width, height, text, size=18, color=DARK_TEXT, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Outfit" if bold and size >= 20 else "Inter"
    p.alignment = align
    return box

def bullets(slide, left, top, width, height, items, size=15, color=DARK_TEXT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Inter"
        p.space_after = Pt(6)
    return box

def card(slide, left, top, width, height, color=CARD_BG):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def teal_card(slide, left, top, width, height):
    return card(slide, left, top, width, height, TEAL)

def notes(slide, text):
    ns = slide.notes_slide
    tf = ns.notes_text_frame
    if tf is not None:
        tf.text = text
    else:
        # Fallback: add to first placeholder
        for ph in ns.placeholders:
            if ph.placeholder_format.idx == 1:
                ph.text = text
                break

def table(slide, left, top, width, height, data, col_widths=None):
    rows, cols = len(data), len(data[0])
    ts = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height))
    t = ts.table
    if col_widths:
        for i, w in enumerate(col_widths):
            t.columns[i].width = Inches(w)
    for r in range(rows):
        for c in range(cols):
            cell = t.cell(r, c)
            cell.text = data[r][c]
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.name = "Inter"
                if r == 0:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = TEAL
                else:
                    p.font.color.rgb = DARK_TEXT
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = CARD_BG
    return t

def bella(slide, left=10.5, top=3.5, height=3):
    if os.path.exists(BELLA):
        slide.shapes.add_picture(BELLA, Inches(left), Inches(top), height=Inches(height))


# ═══════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(LO_TITLE)

# Use the template's title placeholder
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 0:
        ph.text = "Welcome to the Board"
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(44)
            p.font.bold = True
            p.font.name = "Outfit"

txt(slide, 0.8, 3.8, 8, 0.6, "Your Role. Your Impact. Your Compensation.", 24, MUTED_TEXT)
txt(slide, 0.8, 4.6, 8, 0.4, "501(c)(3) Nonprofit  |  EIN: 42-2912572  |  Greensboro, NC", 13, MUTED_TEXT)
txt(slide, 0.8, 5.2, 8, 0.4, "Presented by Bella, Chief Inspiration Officer", 15, GOLD, True)
bella(slide, 9.5, 2.8, 3.5)

notes(slide, """Welcome everyone. Thank you for being here and for believing in this vision.

I founded Think! Ventures to solve a problem I've seen my entire career -- talented people with great business ideas who never launch because they can't afford a website, a brand, or a business plan.

We change that. We build complete businesses for underserved entrepreneurs in 2 to 5 days, at zero cost to them. And today, I want to show you exactly how this works, what your role is, and how we ALL benefit from building this together.

Let's get into it.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 2: THE PROBLEM
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(LO_BLANK)

txt(slide, 0.8, 0.4, 5, 0.4, "THE PROBLEM", 14, GOLD, True)
txt(slide, 0.8, 0.9, 10, 0.6, "72% of entrepreneurs never launch", 38, TEAL, True)

table(slide, 0.8, 2.0, 11.5, 4.2, [
    ["Barrier", "Traditional Cost", "Our Solution"],
    ["Professional Website", "$5,000 - $25,000", "Built in 1-2 days, $0"],
    ["Brand Identity", "$3,000 - $15,000", "Full brand package, $0"],
    ["Business Plan", "$2,000 - $10,000", "AI-researched plan, $0"],
    ["E-Commerce Setup", "$1,000 - $3,000", "Same-day Stripe store, $0"],
    ["LLC Formation", "$500 - $1,500", "Guided filing, $60 state fee"],
    ["TOTAL", "$11,500 - $54,500", "$0 to the entrepreneur"],
], col_widths=[4.0, 3.5, 4.0])

notes(slide, """Here's the reality. Starting a business costs between 11 thousand and 54 thousand dollars before you even make a single sale.

For most people in our community, that's an impossible barrier. And so 72 percent of aspiring entrepreneurs never launch. Not because they don't have good ideas -- because they don't have the capital.

We eliminate every single one of these costs. We build the entire digital business infrastructure in 2 to 5 days, at zero cost to the entrepreneur.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 3: PROOF - ARLAN LLC
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(LO_BLANK)

txt(slide, 0.8, 0.4, 5, 0.4, "PROOF OF CONCEPT  |  1 of 3", 14, GOLD, True)
txt(slide, 0.8, 0.9, 10, 0.6, "Arlan LLC -- Built in 1 Day", 36, TEAL, True)
txt(slide, 0.8, 1.6, 10, 0.3, "Premium Exterior Lighting & Home Services  |  Greensboro, NC", 13, MUTED_TEXT)

table(slide, 0.8, 2.2, 11.5, 3.8, [
    ["Deliverable", "Status", "Market Value"],
    ["Market research & business plan", "Complete", "$5,000"],
    ["Professional brand identity", "Complete", "$8,000"],
    ["8-page premium website", "Complete", "$15,000"],
    ["E-commerce merch store (8 products)", "Complete", "$3,000"],
    ["Stripe payment integration", "Complete", "$2,000"],
    ["AI-powered assistant with TTS", "Complete", "$5,000"],
    ["TOTAL VALUE CREATED", "1 DAY", "$41,500"],
], col_widths=[5.0, 2.5, 4.0])

notes(slide, """Our first proof of concept. Arlan LLC -- premium exterior lighting and home services out of Greensboro.

We built their entire business from scratch in ONE day. Full market research. Professional brand. 8-page premium website. E-commerce store with 8 products. Live Stripe payments. AI-powered customer assistant.

Total market value: 41 thousand 500 dollars. In one day. At zero cost to the entrepreneur.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 4: PROOF - HOOD HYMNS PUBLISHING
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(LO_BLANK)

txt(slide, 0.8, 0.4, 5, 0.4, "PROOF OF CONCEPT  |  2 of 3", 14, GOLD, True)
txt(slide, 0.8, 0.9, 10, 0.6, "Hood Hymns Publishing -- Built in 2 Days", 36, TEAL, True)
txt(slide, 0.8, 1.6, 10, 0.3, "DJ VanHook  |  Literary Publishing Platform  |  hoodhymnspublishing.com", 13, MUTED_TEXT)

table(slide, 0.8, 2.2, 11.5, 3.8, [
    ["Deliverable", "Status", "Market Value"],
    ["Full Next.js publishing platform", "Complete", "$20,000"],
    ["Trilingual i18n (EN/ES/ZH)", "Complete", "$5,000"],
    ["AI narration with Gemini TTS", "Complete", "$8,000"],
    ["Cinematic book trailers", "Complete", "$3,000"],
    ["Merch store with Stripe", "Complete", "$3,000"],
    ["Author dashboard & admin", "Complete", "$5,000"],
    ["TOTAL VALUE CREATED", "2 DAYS", "$44,000"],
], col_widths=[5.0, 2.5, 4.0])

notes(slide, """Proof of concept number two. Hood Hymns Publishing -- a full literary publishing platform for author DJ VanHook.

This one was even more complex. A Next.js application with trilingual support in English, Spanish, and Mandarin. AI-powered narration using Gemini text-to-speech. Cinematic book trailers. A merch store with live Stripe payments. An admin dashboard for managing content.

Total market value: 44 thousand dollars. Built in 2 days.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 5: PROOF - LAKE JEANETTE DENTISTRY
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(LO_BLANK)

txt(slide, 0.8, 0.4, 5, 0.4, "PROOF OF CONCEPT  |  3 of 3", 14, GOLD, True)
txt(slide, 0.8, 0.9, 10, 0.6, "Lake Jeanette Dentistry -- Built in 1 Day", 36, TEAL, True)
txt(slide, 0.8, 1.6, 10, 0.3, "Brenes Precision Dentistry  |  Professional Services  |  Greensboro, NC", 13, MUTED_TEXT)

table(slide, 0.8, 2.2, 11.5, 3.8, [
    ["Deliverable", "Status", "Market Value"],
    ["Professional dental practice website", "Complete", "$12,000"],
    ["Brand identity & visual design", "Complete", "$5,000"],
    ["Patient appointment system", "Complete", "$4,000"],
    ["Service pages & team profiles", "Complete", "$3,000"],
    ["Mobile-responsive design", "Complete", "$2,000"],
    ["SEO & local search optimization", "Complete", "$3,000"],
    ["TOTAL VALUE CREATED", "1 DAY", "$29,000"],
], col_widths=[5.0, 2.5, 4.0])

bella(slide, 10.8, 3.5, 2.5)

notes(slide, """Proof of concept number three. Lake Jeanette Dentistry -- Brenes Precision Dentistry. A professional dental practice here in Greensboro.

Full practice website with service pages, team profiles, appointment system, mobile-responsive design, and local SEO optimization.

Total market value: 29 thousand dollars. One day.

Three businesses. Three days total. Over 114 thousand dollars in combined market value. Zero cost to the entrepreneurs. That's what Think! Ventures delivers.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 4: THE MODEL
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(LO_BLANK)

txt(slide, 0.8, 0.4, 5, 0.4, "THE MODEL", 14, GOLD, True)
txt(slide, 0.8, 0.9, 10, 0.6, "How Nonprofit + For-Profit Work Together", 34, TEAL, True)

# Nonprofit box
teal_card(slide, 0.8, 2.0, 5.5, 4.2)
txt(slide, 1.1, 2.2, 5, 0.4, "THINK! VENTURES FOUNDATION", 17, GOLD, True)
txt(slide, 1.1, 2.7, 5, 0.3, "501(c)(3) Nonprofit", 13, WHITE)
bullets(slide, 1.1, 3.1, 5, 3, [
    "Applies for and receives grants",
    "Recruits and selects entrepreneurs",
    "Runs workshops and programs",
    "Measures and reports social impact",
    "Issues tax-deductible donation receipts",
    "Manages cooperative partner agreements"], 13, WHITE)

# Arrow
txt(slide, 6.5, 3.8, 0.8, 0.6, ">>>", 26, GOLD, True, PP_ALIGN.CENTER)

# LLC box
card(slide, 7.0, 2.0, 5.5, 4.2)
txt(slide, 7.3, 2.2, 5, 0.4, "THINK! DESIGN & PLANNING, LLC", 17, EMERALD, True)
txt(slide, 7.3, 2.7, 5, 0.3, "For-Profit Company (Chris Harrison, Owner)", 13, MUTED_TEXT)
bullets(slide, 7.3, 3.1, 5, 3, [
    "Provides technology services at fair market rates",
    "Builds websites, AI tools, e-commerce",
    "Paid BY the Foundation with grant funds",
    "Also serves direct clients independently",
    "Owns proprietary tech (AVA, Gemini tools)",
    "This is how the Founder gets compensated"], 13, DARK_TEXT)

notes(slide, """Here's how the money flows.

On the left, Think! Ventures Foundation -- the nonprofit. Applies for grants, recruits entrepreneurs, runs workshops, measures impact.

On the right, Think! Design and Planning, LLC -- my for-profit company. Builds the actual technology.

The Foundation contracts with the LLC at fair market rates. 100 percent legal, 100 percent standard. Rates must be at or below market value, the relationship is disclosed on Form 990, and the board approves with me recusing.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 5: THE BOARD
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(LO_BLANK)

txt(slide, 0.8, 0.4, 5, 0.4, "YOUR BOARD", 14, GOLD, True)
txt(slide, 0.8, 0.9, 10, 0.6, "The Leadership Team", 38, TEAL, True)

members = [
    ("Chris Harrison", "Chair / Executive Director", "Founder, NC A&T faculty.\nStrategy, partnerships,\nprogram delivery."),
    ("Gregory Fulton", "Vice Chair", "Assists the Chair.\nCommunity outreach\nand recruitment."),
    ("Damon Howell", "Secretary / Treasurer", "Financials, records,\nForm 990 compliance."),
    ("Charlie Hopper", "Director (At-Large)", "Fresh perspective,\ncommunity voice."),
    ("TBD", "Director (At-Large)", "Independent voice.\nTo be recruited."),
]

for i, (name, role, desc) in enumerate(members):
    x = 0.5 + i * 2.55
    card(slide, x, 2.0, 2.35, 4.2)
    
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.65), Inches(2.2), Inches(1), Inches(1))
    circle.fill.solid()
    circle.fill.fore_color.rgb = CARD_BG
    circle.line.color.rgb = GOLD
    circle.line.width = Pt(2)
    txt(slide, x + 0.65, 2.4, 1, 0.5, "PHOTO", 10, MUTED_TEXT, False, PP_ALIGN.CENTER)
    
    txt(slide, x + 0.1, 3.4, 2.15, 0.4, name, 15, TEAL, True, PP_ALIGN.CENTER)
    txt(slide, x + 0.1, 3.8, 2.15, 0.3, role, 11, GOLD, True, PP_ALIGN.CENTER)
    txt(slide, x + 0.1, 4.2, 2.15, 1.5, desc, 10, MUTED_TEXT, False, PP_ALIGN.CENTER)

notes(slide, """This is our leadership team.

Chris Harrison -- Chair and Executive Director. Strategy, partnerships, program delivery.

Gregory Fulton -- Vice Chair. Community outreach and partner recruitment.

Damon Howell -- Secretary-Treasurer. Financial oversight, records, Form 990.

Charlie Hopper -- Director at Large. Independent community voice.

We have one more seat to fill. Having 5 board members lets all three founders receive compensation while maintaining an independent majority.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 6: WHAT BOARD MEMBERS DO
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(LO_BLANK)

txt(slide, 0.8, 0.4, 5, 0.4, "YOUR COMMITMENT", 14, GOLD, True)
txt(slide, 0.8, 0.9, 10, 0.6, "What Board Members Actually Do", 38, TEAL, True)

# Left - Board duties
card(slide, 0.8, 2.0, 5.5, 4.2)
txt(slide, 1.1, 2.2, 5, 0.4, "Board Service (Unpaid)", 19, EMERALD, True)
bullets(slide, 1.1, 2.8, 5, 3.2, [
    "4 quarterly meetings per year (~1 hour each)",
    "Vote on budgets and major decisions",
    "Review financial reports",
    "Approve contracts and partnerships",
    "Sign annual conflict-of-interest forms",
    "Total time: ~8-10 hours per YEAR"], 14, DARK_TEXT)

# Right - Operational
card(slide, 7.0, 2.0, 5.5, 4.2)
txt(slide, 7.3, 2.2, 5, 0.4, "Operational Work (PAID)", 19, GOLD, True)
bullets(slide, 7.3, 2.8, 5, 3.2, [
    "Separate contractor role for real work",
    "Community outreach & recruitment",
    "Workshop facilitation",
    "Financial operations & grant reporting",
    "Compensated at fair market rates",
    "Approved by the other board members"], 14, DARK_TEXT)

notes(slide, """There are TWO hats you wear.

Board Service -- 4 meetings a year, about an hour each. You vote on budgets, review financials, approve contracts. That's unpaid.

Operational Work -- this is where you get paid. Community outreach, workshop facilitation, financial operations. Compensated as a contractor at fair market rates.

The bylaws specifically allow this. Each person recuses from the vote on their own compensation.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 7: HOW EVERYONE EATS
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(LO_BLANK)

txt(slide, 0.8, 0.4, 5, 0.4, "COMPENSATION", 14, GOLD, True)
txt(slide, 0.8, 0.9, 10, 0.6, "How Everyone Eats", 38, TEAL, True)

# Left card - Board service (everyone)
card(slide, 0.8, 2.0, 5.5, 4.2)
txt(slide, 1.1, 2.2, 5, 0.4, "All Board Members", 19, EMERALD, True)
bullets(slide, 1.1, 2.8, 5, 3.2, [
    "Meeting stipends: $100 - $250 per meeting",
    "4 meetings per year",
    "Approved in annual budget",
    "",
    "Every seat at this table has value.",
    "Every voice gets compensated for showing up."], 14, DARK_TEXT)

# Right card - Operational roles (as work is performed)
card(slide, 7.0, 2.0, 5.5, 4.2)
txt(slide, 7.3, 2.2, 5, 0.4, "Operational Roles (As Needed)", 19, GOLD, True)
bullets(slide, 7.3, 2.8, 5, 3.2, [
    "Executive Director: part-time to full-time salary",
    "Technology services: LLC contract at market rates",
    "Outreach & recruitment: $20 - $35/hr",
    "Financial operations: $20 - $35/hr",
    "Workshop facilitation: $20 - $35/hr",
    "Any member can take on operational work."], 14, DARK_TEXT)

notes(slide, """Here's how compensation works. And I want to be clear -- everyone at this table eats.

On the left -- Board Service. Every board member receives meeting stipends, 100 to 250 dollars per meeting, 4 meetings a year. That's for showing up, voting, reviewing financials. Every seat at this table has value and every voice gets compensated for being here.

On the right -- Operational Roles. These are separate from board duties. The Executive Director draws a salary that grows with grants. My LLC gets paid for technology builds at fair market rates. And anyone on the board who takes on operational work -- outreach, financial ops, workshop facilitation -- gets compensated at 20 to 35 dollars an hour.

The key is this: board service is one hat. Operational work is a second hat. You can wear both. The board votes on each person's operational compensation with that person recused. Everything is transparent, everything is approved, everything is on the books.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 8: GROWTH TIMELINE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(LO_BLANK)

txt(slide, 0.8, 0.4, 5, 0.4, "THE TIMELINE", 14, GOLD, True)
txt(slide, 0.8, 0.9, 10, 0.6, "From Startup to Sustainability", 38, TEAL, True)

table(slide, 0.8, 2.0, 11.5, 4.2, [
    ["", "Phase 1: Pre-Grant", "Phase 2: First Grants", "Phase 3: Growth"],
    ["Foundation Revenue", "$0", "$58K - $165K", "$285K - $750K"],
    ["Executive Director", "LLC income only", "$30K - $60K/yr", "$80K - $150K/yr"],
    ["Operational Contractors", "Volunteer (log hours)", "$15K - $35K/yr each", "$40K - $55K/yr each"],
    ["Board Stipends", "Waived", "$400 - $1,000/yr each", "$1,000 - $2,500/yr each"],
    ["Businesses Launched", "3 (proven)", "10 - 20", "50+"],
    ["Timeline", "Now", "Months 3 - 12", "Year 2 - 3"],
], col_widths=[2.8, 2.8, 3.0, 3.0])

notes(slide, """Three phases.

Phase 1 is right now. Pre-grant. The Executive Director lives off LLC income. Operational team members volunteer and log their hours. Board stipends are waived until we're funded. But we're already building -- 3 businesses launched as proof.

Phase 2 -- first grants hit. Foundation revenue of 58 to 165 thousand. The Executive Director draws 30 to 60 thousand. Operational contractors earn 15 to 35 thousand each. Board stipends kick in at 400 to 1,000 per year per member. We're launching 10 to 20 businesses.

Phase 3 -- growth. 285 to 750 thousand in revenue. Executive Director at 80 to 150 thousand. Operational contractors at 40 to 55 thousand each. Full board stipends. 50 or more businesses launched per year.

The key point -- log your hours now. When the grants arrive, those logged hours establish your track record and justify your compensation going forward.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 9: FUNDING SOURCES
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(LO_BLANK)

txt(slide, 0.8, 0.4, 5, 0.4, "FUNDING STRATEGY", 14, GOLD, True)
txt(slide, 0.8, 0.9, 10, 0.6, "Where the Money Comes From", 38, TEAL, True)

table(slide, 0.5, 2.0, 12.3, 4.4, [
    ["Revenue Source", "Year 1", "Year 2", "Year 3"],
    ["Federal Grants (EDA, SBA, MBDA)", "$25K-$75K", "$50K-$150K", "$100K-$300K"],
    ["State Grants (NC IDEA, NC Commerce)", "$10K-$25K", "$25K-$50K", "$50K-$100K"],
    ["Foundation Grants (Kauffman, Google)", "$5K-$15K", "$15K-$50K", "$25K-$75K"],
    ["HBCU Grants (HBCUFI, Black Ambition)", "$10K-$25K", "$20K-$50K", "$30K-$75K"],
    ["Corporate Sponsorships", "$0", "$5K-$15K", "$10K-$30K"],
    ["Donations + Workshop Fees", "$3K-$10K", "$10K-$30K", "$20K-$50K"],
    ["10% Cooperative Profit Share", "$5K-$15K", "$20K-$50K", "$50K-$120K"],
    ["TOTAL", "$58K-$165K", "$145K-$395K", "$285K-$750K"],
], col_widths=[4.0, 2.5, 2.8, 3.0])

notes(slide, """The biggest source is grants. Federal, state, foundation, and HBCU-specific. We've identified over 15 high-fit opportunities.

We also generate revenue through the cooperative profit share. Every business shares 10 percent of profits.

By Year 3, total budget of 285 to 750 thousand dollars.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 10: THE ASK
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(LO_BLANK)

txt(slide, 0.8, 0.4, 5, 0.4, "THE ASK", 14, GOLD, True)
txt(slide, 0.8, 0.9, 10, 0.6, "One Simple Step", 38, TEAL, True)

card(slide, 1.5, 2.2, 8, 3.8)
txt(slide, 2.0, 2.5, 7, 0.6, "If you believe in what you just saw:", 22, MUTED_TEXT)
txt(slide, 2.0, 3.2, 7, 0.8, "Sign the bylaws.", 36, TEAL, True)
txt(slide, 2.0, 4.2, 7, 1.2,
    "That's it. I handle everything else --\nthe filings, the grants, the tech, the builds.\nYou showed up. That's what matters.",
    18, DARK_TEXT)

bella(slide, 10.0, 2.8, 3)

notes(slide, """Here's the thing. I'm not asking you to do a hundred things.

If you believe in what you just saw -- the model, the proof, the vision -- then all I need is your signature on the bylaws.

I handle everything else. The IRS filing. The grant applications. The technology builds. The bank account. All of it.

You showed up. You believed. That's what matters. Sign the bylaws and let's get to work.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 11: THE VISION
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(LO_SECTION)

txt(slide, 1.5, 1.5, 10, 0.6, "THE VISION", 16, GOLD, True, PP_ALIGN.CENTER)

txt(slide, 1.5, 2.5, 10, 1.5,
    "\"A world where lack of capital, technical knowledge,\nor professional connections never prevents a great idea\nfrom becoming a thriving business.\"",
    26, TEAL, False, PP_ALIGN.CENTER)

txt(slide, 1.5, 4.5, 10, 0.5, "From Dream to Launch in Days.", 30, GOLD, True, PP_ALIGN.CENTER)
txt(slide, 1.5, 5.5, 10, 0.4, "thinkventures.app", 16, MUTED_TEXT, False, PP_ALIGN.CENTER)

bella(slide, 5.5, 5.0, 2)

notes(slide, """Our vision is a world where lack of capital, technical knowledge, or professional connections never prevents a great idea from becoming a thriving business.

From dream to launch in days. We proved it with Arlan. We'll prove it 50 more times, then 500 more times.

You're not just joining a board. You're joining a movement. Let's build something extraordinary together.

Now -- let's sign these documents and get to work.""")

# ═══════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════
prs.save(OUTPUT)
print(f"Presentation saved to: {OUTPUT}")
print(f"Total slides: {len(prs.slides)}")
print("Built on Think! Ventures official template with all branding baked in.")
