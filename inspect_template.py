"""Inspect the Think Ventures .potx template"""
from zipfile import ZipFile
from pptx import Presentation
import shutil

src = r'C:\Users\Chris\Desktop\WEBSITES\THINK! VENTURES\docs\Think_Ventures_Template.potx'
dst = r'C:\Users\Chris\Desktop\WEBSITES\THINK! VENTURES\docs\Think_Ventures_Template_converted.pptx'

# Convert .potx to .pptx by changing the content type
shutil.copy2(src, dst)

with ZipFile(dst, 'r') as zin:
    all_data = {name: zin.read(name) for name in zin.namelist()}

ct = all_data['[Content_Types].xml'].decode('utf-8')
ct = ct.replace(
    'application/vnd.openxmlformats-officedocument.presentationml.template.main+xml',
    'application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml'
)
all_data['[Content_Types].xml'] = ct.encode('utf-8')

with ZipFile(dst, 'w') as zout:
    for name, data in all_data.items():
        zout.writestr(name, data)

prs = Presentation(dst)
print(f"Slide width: {prs.slide_width / 914400:.2f} in")
print(f"Slide height: {prs.slide_height / 914400:.2f} in")
print(f"Slide layouts: {len(prs.slide_layouts)}")
print(f"Existing slides: {len(prs.slides)}")
print()
for i, layout in enumerate(prs.slide_layouts):
    print(f'  Layout {i}: "{layout.name}"')
    for ph in layout.placeholders:
        print(f'    PH {ph.placeholder_format.idx}: {ph.name} ({ph.placeholder_format.type})')
print()
print("Template converted and ready at:", dst)
