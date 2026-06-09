"""
Think! Ventures Foundation -- Arlan Franchise & Partnership Meeting
Presentation for Dylan Thomas & Partner
Built on the official Think! Ventures .potx template
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from zipfile import ZipFile
import shutil, os

# ─── PATHS ───
TEMPLATE_SRC = r'C:\Users\Chris\Desktop\WEBSITES\THINK! VENTURES\docs\Think_Ventures_Template.potx'
TEMPLATE      = r'C:\Users\Chris\Desktop\WEBSITES\THINK! VENTURES\docs\Think_Ventures_Template_converted.pptx'
OUTPUT        = r'C:\Users\Chris\Downloads\Arlan_Partnership_Meeting.pptx'
BELLA         = r'C:\Users\Chris\Desktop\WEBSITES\THINK! VENTURES\assets\images\bella-mascot.png'

# ─── BRAND COLORS ───
TEAL       = RGBColor(0x0D, 0x4F, 0x4F)
GOLD       = RGBColor(0xF5, 0xA6, 0x23)
EMERALD    = RGBColor(0x10, 0xB9, 0x81)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT  = RGBColor(0x1A, 0x2A, 0x3A)
MUTED_TEXT = RGBColor(0x4A, 0x5A, 0x6A)
CARD_BG    = RGBColor(0xE8, 0xED, 0xF2)
LIGHT_GOLD = RGBColor(0xFD, 0xF0, 0xD5)
RED        = RGBColor(0xEF, 0x44, 0x44)

# ─── CONVERT .potx TO .pptx ───
if not os.path.exists(TEMPLATE):
    shutil.copy2(TEMPLATE_SRC, TEMPLATE)
    with ZipFile(TEMPLATE, 'r') as zin:
        all_data = {n: zin.read(n) for n in zin.namelist()}
    ct = all_data['[Content_Types].xml'].decode('utf-8').replace(
        'presentationml.template.main+xml',
        'presentationml.presentation.main+xml')
    all_data['[Content_Types].xml'] = ct.encode('utf-8')
    with ZipFile(TEMPLATE, 'w') as zout:
        for n, d in all_data.items():
            zout.writestr(n, d)

prs = Presentation(TEMPLATE)

# Layout references
LO_TITLE   = prs.slide_layouts[0]  # Title Slide
LO_CONTENT = prs.slide_layouts[1]  # Title and Content
LO_SECTION = prs.slide_layouts[2]  # Section Header
LO_BLANK   = prs.slide_layouts[3]  # Blank

# ─── HELPERS ───
def txt(slide, left, top, width, height, text, size=18, color=DARK_TEXT, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Outfit" if bold and size >= 20 else "Inter"
    p.alignment = align
    return box

def bullets(slide, left, top, width, height, items, size=15, color=DARK_TEXT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Inter"
        p.space_after = Pt(6)
    return box

def card(slide, left, top, width, height, color=CARD_BG):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def teal_card(slide, left, top, width, height):
    return card(slide, left, top, width, height, TEAL)

def gold_card(slide, left, top, width, height):
    return card(slide, left, top, width, height, LIGHT_GOLD)

def notes(slide, text):
    ns = slide.notes_slide
    tf = ns.notes_text_frame
    if tf is not None:
        tf.text = text
    else:
        for ph in ns.placeholders:
            if ph.placeholder_format.idx == 1:
                ph.text = text
                break

def table(slide, left, top, width, height, data, col_widths=None):
    rows, cols = len(data), len(data[0])
    ts = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height))
    t = ts.table
    if col_widths:
        for i, w in enumerate(col_widths):
            t.columns[i].width = Inches(w)
    for r in range(rows):
        for c in range(cols):
            cell = t.cell(r, c)
            cell.text = data[r][c]
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.name = "Inter"
                if r == 0:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = TEAL
                else:
                    p.font.color.rgb = DARK_TEXT
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = CARD_BG
    return t

def bella(slide, left=10.5, top=3.5, height=3):
    if os.path.exists(BELLA):
        slide.shapes.add_picture(BELLA, Inches(left), Inches(top), height=Inches(height))


# ═══════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(LO_TITLE)

for ph in slide.placeholders:
    if ph.placeholder_format.idx == 0:
        ph.text = "Arlan LLC × Think! Ventures"
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(40)
            p.font.bold = True
            p.font.name = "Outfit"

txt(slide, 0.8, 3.8, 8, 0.6, "Your Business. Our Infrastructure. Let's Build.", 24, MUTED_TEXT)
txt(slide, 0.8, 4.6, 8, 0.4, "Partnership Meeting  |  June 8, 2026  |  Greensboro, NC", 13, MUTED_TEXT)
txt(slide, 0.8, 5.2, 8, 0.4, "Presented by Chris Harrison — Think! Ventures Foundation", 15, GOLD, True)
bella(slide, 9.5, 2.8, 3.5)

notes(slide, """Welcome, Dylan and [friend]. Thank you for taking this meeting.

Today I'm going to show you three things:

1. What we've already built for Arlan LLC — the value on the table right now
2. How Think! Ventures works and why this is different from anything else
3. The franchise model — what this looks like monetarily for you

Let's get into it.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 2: WHAT WE ALREADY BUILT
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(LO_BLANK)

txt(slide, 0.8, 0.4, 5, 0.4, "WHAT YOU HAVE RIGHT NOW", 14, GOLD, True)
txt(slide, 0.8, 0.9, 10, 0.6, "Arlan LLC — Already Built", 38, TEAL, True)
txt(slide, 0.8, 1.6, 10, 0.3, "arlanpro.com  |  Premium Exterior Lighting & Home Services", 13, MUTED_TEXT)

table(slide, 0.8, 2.2, 11.5, 3.8, [
    ["Deliverable", "Status", "Market Value"],
    ["Market research & business plan", "✓ Complete", "$5,000"],
    ["Professional brand identity (logo, mascot, colors)", "✓ Complete", "$8,000"],
    ["8-page premium website with animations", "✓ Complete", "$15,000"],
    ["E-commerce merch store (8 products)", "✓ Complete", "$3,000"],
    ["Stripe payment processing (live)", "✓ Complete", "$2,000"],
    ["AI-powered assistant with voice (Arlan)", "✓ Complete", "$5,000"],
    ["Drone 3D experience page", "✓ Complete", "$3,500"],
    ["TOTAL VALUE DELIVERED", "1 DAY", "$41,500"],
], col_widths=[5.0, 2.5, 4.0])

notes(slide, """This is what's already been built for Arlan LLC. All of this is live right now at arlanpro.com.

Full market research and business plan. Professional brand identity -- logo, mascot, color palette. An 8-page premium website with smooth animations. An e-commerce merch store with 8 products. Live Stripe payment processing. An AI-powered customer assistant with text-to-speech voice. And a drone 3D experience page.

Total market value: 41 thousand 500 dollars. Built in ONE day. Cost to Dylan: zero dollars.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 3: WHAT IS THINK! VENTURES
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(LO_BLANK)

txt(slide, 0.8, 0.4, 5, 0.4, "THE MODEL", 14, GOLD, True)
txt(slide, 0.8, 0.9, 10, 0.6, "What is Think! Ventures?", 38, TEAL, True)

# Left card
teal_card(slide, 0.8, 2.0, 5.5, 4.5)
txt(slide, 1.1, 2.2, 5, 0.4, "THINK! VENTURES FOUNDATION", 17, GOLD, True)
txt(slide, 1.1, 2.7, 5, 0.3, "501(c)(3) Nonprofit", 13, WHITE)
bullets(slide, 1.1, 3.1, 5, 3.2, [
    "Funded by grants (federal, state, foundation)",
    "Builds entire digital businesses for free",
    "Brand, website, merch, business plan — everything",
    "Partners with entrepreneurs who can't afford agencies",
    "You bring the dream. We build the infrastructure.",
    "From dream to launch in DAYS, not months"], 13, WHITE)

# Arrow
txt(slide, 6.5, 3.8, 0.8, 0.6, ">>>", 26, GOLD, True, PP_ALIGN.CENTER)

# Right card
card(slide, 7.0, 2.0, 5.5, 4.5)
txt(slide, 7.3, 2.2, 5, 0.4, "THE DEAL", 17, EMERALD, True)
txt(slide, 7.3, 2.7, 5, 0.3, "Simple Cooperative Partnership", 13, MUTED_TEXT)
bullets(slide, 7.3, 3.1, 5, 3.2, [
    "You pay: $0 upfront",
    "You pay: $0 monthly",
    "You pay: $0 for the website, brand, or tools",
    "",
    "When you profit, 10% goes back to Think!",
    "When you DON'T profit, you owe NOTHING"], 13, DARK_TEXT)

notes(slide, """Think! Ventures Foundation is a 501(c)(3) nonprofit. We're funded by grants — federal, state, and foundation grants.

What do we do? We build entire digital businesses for free. Brand, website, merch store, business plan, AI assistant — everything you need to launch. For zero dollars.

The deal is simple. You pay nothing upfront. Nothing monthly. Nothing for the website, brand, or tools. When your business becomes profitable, 10 percent goes back to Think! Ventures to fund the next entrepreneur. When you don't make money, you owe us nothing. We succeed together or not at all.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 4: WHY THIS IS DIFFERENT
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(LO_BLANK)

txt(slide, 0.8, 0.4, 5, 0.4, "WHY THIS IS DIFFERENT", 14, GOLD, True)
txt(slide, 0.8, 0.9, 10, 0.6, "Traditional vs. Think! Ventures", 36, TEAL, True)

table(slide, 0.8, 2.0, 11.5, 4.5, [
    ["", "Traditional Route", "Think! Ventures"],
    ["Startup Cost", "$15,000 - $55,000", "$0"],
    ["Timeline", "3 - 6 months", "2 - 5 days"],
    ["Website", "Template from Wix/Squarespace", "Custom premium build"],
    ["Brand", "DIY Canva logo", "Professional identity package"],
    ["Merch Store", "Not included", "8+ products, live payments"],
    ["Business Plan", "$2K from consultant", "AI-researched, bank-ready"],
    ["AI Assistant", "Not available", "Custom voice bot on your site"],
    ["Ongoing Support", "Pay per hour", "Included in partnership"],
], col_widths=[2.8, 4.0, 4.7])

notes(slide, """Let me show you why this is completely different from anything else out there.

The traditional route costs 15 to 55 thousand dollars and takes 3 to 6 months. You get a template website from Wix, a DIY Canva logo, and no merch, no business plan, no AI.

With Think! Ventures, you get all of that in 2 to 5 days for zero dollars. Custom premium website, professional brand, merch store with live payments, AI-researched business plan, custom AI assistant, and ongoing support.

This is not a template. This is not Wix. This is a 41-thousand-dollar digital business built by AI-powered technology in a single day.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 5: THE FRANCHISE CONCEPT
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(LO_BLANK)

txt(slide, 0.8, 0.4, 5, 0.4, "THE FRANCHISE OPPORTUNITY", 14, GOLD, True)
txt(slide, 0.8, 0.9, 10, 0.6, "Arlan as a Repeatable Model", 38, TEAL, True)

# The idea
card(slide, 0.8, 2.0, 11.5, 1.8)
txt(slide, 1.3, 2.2, 10, 0.6, "The Big Idea", 22, TEAL, True)
txt(slide, 1.3, 2.8, 10, 0.7,
    "Arlan's model — premium exterior lighting & home services — works in ANY market.\n"
    "We build the brand, website, and systems ONCE, then replicate in multiple cities.", 15, DARK_TEXT)

# How it works
card(slide, 0.8, 4.2, 3.6, 2.5)
txt(slide, 1.1, 4.4, 3, 0.4, "1. Pick a Market", 17, EMERALD, True)
txt(slide, 1.1, 4.9, 3, 1.5, "Choose a city — Asheville,\nBoone, Tampa, Charlotte.\nResearch the demand.\nCustomize the messaging.", 13, DARK_TEXT)

card(slide, 4.7, 4.2, 3.6, 2.5)
txt(slide, 5.0, 4.4, 3, 0.4, "2. Clone the System", 17, EMERALD, True)
txt(slide, 5.0, 4.9, 3, 1.5, "Duplicate the website.\nLocalize for the new market.\nSet up new Stripe account.\nLaunch merch line.", 13, DARK_TEXT)

card(slide, 8.6, 4.2, 3.6, 2.5)
txt(slide, 8.9, 4.4, 3, 0.4, "3. Operate & Grow", 17, EMERALD, True)
txt(slide, 8.9, 4.9, 3, 1.5, "Local operator runs the biz.\nThink! handles all tech.\nRevenue splits per agreement.\nScale to next market.", 13, DARK_TEXT)

notes(slide, """Here's where it gets exciting. The franchise opportunity.

Arlan's model works in ANY market. Premium exterior lighting and home services -- holiday lights, permanent LED, landscape lighting, window cleaning, drone inspections. Every city needs this.

Step 1: Pick a market. Asheville, Boone, Tampa, Charlotte -- wherever there's demand.

Step 2: Clone the system. We duplicate the website, localize it for the new market, set up a new Stripe account, launch a local merch line. This takes us ONE day because the template already exists.

Step 3: Operate and grow. A local operator runs the business day-to-day. Think! handles all the technology. Revenue splits per agreement. Then we scale to the next market.

This is the franchise without the franchise fee. You're not paying 50 thousand dollars to buy into a franchise. You're getting the entire system for free.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 6: FINANCIAL PROJECTIONS — SINGLE MARKET
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(LO_BLANK)

txt(slide, 0.8, 0.4, 5, 0.4, "FINANCIAL PROJECTIONS", 14, GOLD, True)
txt(slide, 0.8, 0.9, 10, 0.6, "What One Market Looks Like", 38, TEAL, True)
txt(slide, 0.8, 1.6, 10, 0.3, "Based on Arlan LLC — Premium Exterior Lighting & Home Services", 13, MUTED_TEXT)

table(slide, 0.5, 2.2, 12.3, 4.2, [
    ["Revenue Stream", "Year 1 (Startup)", "Year 2 (Growth)", "Year 3 (Established)"],
    ["Holiday Lighting (10→25→40 installs)", "$25,000", "$62,500", "$100,000"],
    ["Permanent LED (5→15→25 installs)", "$17,500", "$52,500", "$87,500"],
    ["Landscape Lighting (5→12→20 jobs)", "$7,500", "$18,000", "$30,000"],
    ["Window Cleaning (monthly clients)", "$6,000", "$18,000", "$36,000"],
    ["Drone Inspections ($150/ea)", "$3,000", "$7,500", "$15,000"],
    ["Merch Sales (online + events)", "$2,000", "$5,000", "$10,000"],
    ["TOTAL GROSS REVENUE", "$61,000", "$163,500", "$278,500"],
    ["Operator Take-Home (after expenses)", "$30,500", "$81,750", "$153,175"],
], col_widths=[3.5, 2.8, 3.0, 3.0])

notes(slide, """Let's talk real numbers. One market. One operator.

Year 1 is startup. You're building the client base. 10 holiday installs, 5 permanent LED jobs, some landscape lighting, window cleaning, drone inspections, and merch sales. Total gross revenue: 61 thousand dollars. After expenses, your take-home is about 30 thousand.

Year 2 you're growing. Word of mouth kicks in, your portfolio is building, referrals are coming. 163 thousand gross, 81 thousand take-home.

Year 3 you're established. You're the go-to lighting company in your market. 278 thousand gross, 153 thousand take-home.

These numbers are conservative. The margins on holiday lighting are 40 to 60 percent. Permanent LED installations are even higher because they're one-time installs with recurring app-control revenue.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 7: THE MONEY SPLIT
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(LO_BLANK)

txt(slide, 0.8, 0.4, 5, 0.4, "THE MONEY SPLIT", 14, GOLD, True)
txt(slide, 0.8, 0.9, 10, 0.6, "How Everyone Eats", 38, TEAL, True)

# Visual money flow
gold_card(slide, 0.8, 2.0, 11.5, 1.5)
txt(slide, 1.3, 2.1, 5, 0.5, "Example: Year 2 Revenue = $163,500", 22, TEAL, True)
txt(slide, 1.3, 2.7, 10, 0.5, "Here's exactly where every dollar goes -- complete transparency.", 14, MUTED_TEXT)

table(slide, 0.8, 3.8, 11.5, 3.2, [
    ["Expense", "% of Revenue", "Annual Amount", "Who Gets It"],
    ["Materials & Supplies", "25%", "$40,875", "Suppliers"],
    ["Operating Expenses (insurance, truck, tools)", "15%", "$24,525", "Business costs"],
    ["Co-op Share (10% of NET profit)", "10% of profit", "$9,810", "Co-op Pool (see next slide)"],
    ["YOUR NET TAKE-HOME", "50%", "$81,750", "Dylan's pocket"],
    ["Taxes (estimated)", "~15% of net", "$12,263", "IRS / State"],
    ["ACTUAL CASH IN YOUR POCKET", "", "$69,487", "YOU"],
], col_widths=[3.5, 2.0, 2.5, 3.5])

notes(slide, """Let me break down exactly where the money goes. Year 2 example -- 163 thousand in revenue.

25 percent goes to materials and supplies. 15 percent to operating expenses -- insurance, truck, tools.

That leaves about 98 thousand in NET profit.

10 percent of that net profit -- 9,810 dollars -- goes into the cooperative pool. But here's the thing -- that money doesn't just disappear. Almost half of it comes BACK to you as savings and dividends. I'll show you exactly how on the next slide.

Your take-home before co-op returns: 81 thousand 750 dollars. After taxes, about 69 thousand in your pocket. PLUS the co-op returns.

And remember -- you paid ZERO to get started. No franchise fee. No website cost. No branding cost.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 8: CO-OP DIVIDEND MODEL
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(LO_BLANK)

txt(slide, 0.8, 0.4, 5, 0.4, "THE CO-OP ADVANTAGE", 14, GOLD, True)
txt(slide, 0.8, 0.9, 10, 0.6, "Your 10% Comes Back to You", 38, TEAL, True)
txt(slide, 0.8, 1.6, 10, 0.3, "Like REI, Ace Hardware, and Land O'Lakes -- a cooperative where members share in the success", 13, MUTED_TEXT)

# How the 10% splits - 3 markets example: $29,610 total pool
table(slide, 0.5, 2.2, 12.3, 2.5, [
    ["Co-op Pool Allocation", "% of 10%", "Annual (3 Markets)", "What It Does"],
    ["Think! Foundation (mission)", "40%", "$11,844", "Funds next 2-3 business launches"],
    ["Operator Savings Fund", "20%", "$5,922 (~$1,974 each)", "Your emergency fund / nest egg"],
    ["Annual Patronage Dividends", "20%", "$5,922 (~$1,974 each)", "Cash back at year-end"],
    ["Growth Reinvestment", "20%", "$5,922", "Better tools, group insurance, marketing"],
], col_widths=[3.0, 1.5, 3.5, 4.3])

# Bottom: What Dylan ACTUALLY keeps
gold_card(slide, 0.8, 5.2, 11.5, 1.8)
txt(slide, 1.3, 5.3, 5, 0.5, "What Dylan ACTUALLY Keeps (Year 2)", 20, TEAL, True)

table(slide, 1.3, 5.9, 10.5, 0.9, [
    ["Take-Home Pay", "+ Savings Return", "+ Year-End Dividend", "= REAL TOTAL"],
    ["$81,750", "+ $1,974", "+ $1,974", "$85,698"],
], col_widths=[2.8, 2.5, 2.5, 2.7])

notes(slide, """Here's what makes this different from just paying 10 percent and never seeing it again.

We structure this as a cooperative -- like REI, Ace Hardware, or Land O'Lakes. The 10 percent goes into a shared co-op pool, and almost HALF of it comes back to you.

40 percent of the pool -- about 12 thousand -- funds Think! Ventures Foundation. That's the mission money that launches the next 2-3 businesses.

20 percent -- about 6 thousand -- goes into an Operator Savings Fund. That's YOUR money sitting in an account for emergencies, equipment upgrades, or a rainy day. About 2 thousand per operator per year.

20 percent -- another 6 thousand -- comes back as Annual Patronage Dividends. Cash in your pocket at the end of the year. Like REI's annual member dividend. About 2 thousand per operator.

20 percent -- the last 6 thousand -- goes into Growth Reinvestment. Better tools, group insurance rates, shared marketing. Things that benefit ALL operators.

So what does Dylan actually keep? Not just the 81 thousand take-home -- add back the 2 thousand savings return and the 2 thousand dividend. Your real total is 85 thousand 698 dollars. The 10 percent isn't a fee -- it's an investment in a system that pays you back.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 9: MULTI-MARKET SCALING
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(LO_BLANK)

txt(slide, 0.8, 0.4, 5, 0.4, "SCALING THE FRANCHISE", 14, GOLD, True)
txt(slide, 0.8, 0.9, 10, 0.6, "What 3 Markets Looks Like", 38, TEAL, True)

table(slide, 0.8, 2.0, 11.5, 4.5, [
    ["Market", "Operator", "Year 2 Gross", "Take-Home + Co-op Returns", "Co-op Share"],
    ["Greensboro, NC", "Dylan", "$163,500", "$85,698", "$9,810"],
    ["Asheville / Boone, NC", "Partner (TBD)", "$145,000", "$76,448", "$8,700"],
    ["Tampa Bay, FL", "Partner (TBD)", "$185,000", "$96,448", "$11,100"],
    ["", "", "", "", ""],
    ["COMBINED", "3 operators", "$493,500", "$258,594 total", "$29,610"],
], col_widths=[2.5, 2.0, 2.5, 2.8, 1.7])

txt(slide, 0.8, 6.6, 11, 0.5,
    "Each market costs ~$0 to clone. Every operator earns dividends from the collective success of ALL markets.",
    14, MUTED_TEXT)

bella(slide, 10.5, 4.5, 2.5)

notes(slide, """Now let's scale this. Three markets.

Greensboro with Dylan. Asheville-Boone with a partner. Tampa Bay with a partner.

Year 2 combined: 493 thousand in gross revenue. Operators take home a combined 258 thousand INCLUDING their co-op savings and dividends.

The co-op pool totals about 29 thousand. 12 thousand funds the mission. 12 thousand flows back to operators as savings and dividends. 6 thousand grows everyone's business.

And the more markets we add, the bigger the pool, the bigger the dividends. When there's 10 markets, 20 markets -- every operator benefits from the collective success.

Each market costs Think! essentially zero to clone. We duplicate the site in one day. The operator runs independently. The co-op ties everyone together.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 9: WHAT DYLAN & PARTNER DO
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(LO_BLANK)

txt(slide, 0.8, 0.4, 5, 0.4, "YOUR ROLE", 14, GOLD, True)
txt(slide, 0.8, 0.9, 10, 0.6, "What You Do vs. What We Do", 38, TEAL, True)

# Left - Dylan's role
teal_card(slide, 0.8, 2.0, 5.5, 4.5)
txt(slide, 1.1, 2.2, 5, 0.4, "WHAT YOU DO", 19, GOLD, True)
txt(slide, 1.1, 2.7, 5, 0.3, "Dylan & Partner", 13, WHITE)
bullets(slide, 1.1, 3.1, 5, 3.2, [
    "Sell the service — knock on doors, get clients",
    "Install the lights and do the work",
    "Build relationships with homeowners",
    "Show up to estimates and close deals",
    "Manage your truck, tools, and schedule",
    "Deliver excellent work and earn referrals",
    "",
    "YOU are the face of Arlan."], 13, WHITE)

# Right - Think! handles
card(slide, 7.0, 2.0, 5.5, 4.5)
txt(slide, 7.3, 2.2, 5, 0.4, "WHAT WE HANDLE", 19, EMERALD, True)
txt(slide, 7.3, 2.7, 5, 0.3, "Think! Design & Planning", 13, MUTED_TEXT)
bullets(slide, 7.3, 3.1, 5, 3.2, [
    "Website updates and maintenance",
    "Payment processing (Stripe)",
    "Merch store management",
    "SEO and online presence",
    "AI assistant updates",
    "New market launches (franchise cloning)",
    "All technology and digital infrastructure",
    "WE are the engine behind the scenes."], 13, DARK_TEXT)

notes(slide, """Here's the split.

Dylan and partner -- you are the business. You sell, you install, you build relationships, you deliver excellent work. You are the face of Arlan in your market.

Think! handles everything digital. Website, payments, merch, SEO, AI, new market launches. We are the engine behind the scenes.

You focus on what you're good at -- the craft, the sales, the customer relationships. We focus on what we're good at -- the technology, the digital infrastructure, the scaling.

Neither of us succeeds without the other. That's why it's a cooperative partnership, not just a service agreement.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 10: WHAT YOU NEED TO DO TODAY
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(LO_BLANK)

txt(slide, 0.8, 0.4, 5, 0.4, "NEXT STEPS", 14, GOLD, True)
txt(slide, 0.8, 0.9, 10, 0.6, "What Happens Next", 38, TEAL, True)

# Step cards
steps = [
    ("1", "Form Your LLC", "File with NC Secretary of State ($125).\nGet your EIN (free, instant at IRS.gov).\nOpen a business bank account.\nThink! guides you through every step."),
    ("2", "Sign the Partnership Agreement", "Cooperative agreement between\nArlan LLC and Think! Ventures.\n10% of net profits when profitable.\n$0 when you're not.\nSimple. Fair. Transparent."),
    ("3", "Get Your First Clients", "You already have a $41K website.\nStart knocking on doors.\nBook your first 5 holiday installs.\nWe handle all the online leads.\nGoal: first revenue in 30 days."),
]

for i, (num, title, desc) in enumerate(steps):
    x = 0.8 + i * 4.1
    card(slide, x, 2.0, 3.7, 4.5)
    
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 1.35), Inches(2.2), Inches(0.9), Inches(0.9))
    circle.fill.solid()
    circle.fill.fore_color.rgb = GOLD
    circle.line.fill.background()
    txt(slide, x + 1.35, 2.3, 0.9, 0.7, num, 30, DARK_TEXT, True, PP_ALIGN.CENTER)
    
    txt(slide, x + 0.2, 3.3, 3.2, 0.5, title, 18, TEAL, True, PP_ALIGN.CENTER)
    txt(slide, x + 0.2, 3.9, 3.2, 2.4, desc, 12, DARK_TEXT, False, PP_ALIGN.CENTER)

notes(slide, """Three steps. That's it.

Step 1: Form your LLC. File with North Carolina -- 125 dollars. Get your EIN from the IRS -- free, takes 5 minutes online. Open a business bank account. Think! guides you through every step, and our LaunchPad tool walks you through the exact forms.

Step 2: Sign the cooperative partnership agreement. 10 percent of net profits when you're profitable. Zero when you're not. Simple, fair, transparent.

Step 3: Get your first clients. You already have a 41-thousand-dollar website live at arlanpro.com. Start knocking on doors. Book your first 5 holiday lighting installs. We handle all the online leads that come through the website. Goal: first revenue within 30 days.

You can literally file the LLC today, sign the agreement today, and start selling tomorrow.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 11: THE VISION
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(LO_SECTION)

txt(slide, 1.5, 1.5, 10, 0.6, "THE VISION", 16, GOLD, True, PP_ALIGN.CENTER)

txt(slide, 1.5, 2.5, 10, 1.5,
    "\"You came in with a dream.\nYou're leaving with a $41,500 business.\nNow let's turn it into a $300K empire.\"",
    28, TEAL, False, PP_ALIGN.CENTER)

txt(slide, 1.5, 4.5, 10, 0.5, "From Dream to Launch in Days.", 30, GOLD, True, PP_ALIGN.CENTER)
txt(slide, 1.5, 5.3, 10, 0.4, "arlanpro.com  |  think-ventures.netlify.app", 16, MUTED_TEXT, False, PP_ALIGN.CENTER)

bella(slide, 5.5, 5.0, 2)

notes(slide, """You came in with a dream. You're leaving with a 41 thousand 500 dollar business. Live. Online. Taking payments. Right now.

Now let's turn it into a 300 thousand dollar empire across multiple markets.

Form the LLC. Sign the agreement. Start selling.

From dream to launch in days. That's not a slogan -- you've already lived it. Now let's scale it.

Let's do this.""")

# ═══════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════
prs.save(OUTPUT)
print(f"[OK] Presentation saved to: {OUTPUT}")
print(f"   Total slides: {len(prs.slides)}")
print("   Built on Think! Ventures official template")
print("   Speaker notes included on every slide")
