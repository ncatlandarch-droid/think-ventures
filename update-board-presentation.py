"""
Think! Ventures Foundation -- Board Presentation Updater
Takes Chris's customized "to actually use" PPTX and updates:
 - Board roles (generic executive titles, Cadasia added, CPA placeholder)
 - Compensation slide (role-based, no names)
 - Timeline slide (role-based, no names)
 - Coaching program added to model
 - Speaker notes on all slides
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os, shutil, copy

# ─── PATHS ───
SOURCE = r'C:\Users\Chris\Downloads\Think_Ventures_Board_Final to actually use.pptx'
OUTPUT = r'C:\Users\Chris\Downloads\Think_Ventures_Board_READY.pptx'
BELLA  = r'C:\Users\Chris\Desktop\WEBSITES\THINK! VENTURES\assets\images\bella-mascot.png'

# ─── BRAND COLORS ───
TEAL       = RGBColor(0x0D, 0x4F, 0x4F)
GOLD       = RGBColor(0xF5, 0xA6, 0x23)
EMERALD    = RGBColor(0x10, 0xB9, 0x81)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT  = RGBColor(0x1A, 0x2A, 0x3A)
MUTED_TEXT = RGBColor(0x4A, 0x5A, 0x6A)
CARD_BG    = RGBColor(0xE8, 0xED, 0xF2)

# ─── HELPERS ───

def clear_slide(slide):
    """Remove all shapes from a slide"""
    for shape in list(slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

def txt(slide, left, top, width, height, text, size=18, color=DARK_TEXT, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Outfit" if bold and size >= 18 else "Inter"
    p.alignment = align
    return box

def bullets(slide, left, top, width, height, items, size=14, color=DARK_TEXT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Inter"
        p.space_after = Pt(6)
    return box

def card(slide, left, top, width, height, color=CARD_BG):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def teal_card(slide, left, top, width, height):
    return card(slide, left, top, width, height, TEAL)

def notes(slide, text):
    ns = slide.notes_slide
    tf = ns.notes_text_frame
    if tf is not None:
        tf.text = text
    else:
        for ph in ns.placeholders:
            if ph.placeholder_format.idx == 1:
                ph.text = text
                break

def table(slide, left, top, width, height, data, col_widths=None):
    rows, cols = len(data), len(data[0])
    ts = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height))
    t = ts.table
    if col_widths:
        for i, w in enumerate(col_widths):
            t.columns[i].width = Inches(w)
    for r in range(rows):
        for c in range(cols):
            cell = t.cell(r, c)
            cell.text = data[r][c]
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.name = "Inter"
                if r == 0:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = TEAL
                else:
                    p.font.color.rgb = DARK_TEXT
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = CARD_BG
    return t

def bella(slide, left=10.5, top=3.5, height=3):
    if os.path.exists(BELLA):
        slide.shapes.add_picture(BELLA, Inches(left), Inches(top), height=Inches(height))


# ═══════════════════════════════════════════════════════════════
# LOAD THE PRESENTATION
# ═══════════════════════════════════════════════════════════════
prs = Presentation(SOURCE)
slides = list(prs.slides)
print(f"Loaded: {len(slides)} slides from source")

# ═══════════════════════════════════════════════════════════════
# SLIDE 1: TITLE - just add notes
# ═══════════════════════════════════════════════════════════════
notes(slides[0], """Welcome everyone. Thank you for being here and for believing in this vision.

I founded Think! Ventures to solve a problem I've seen my entire career -- talented people with great business ideas who never launch because they can't afford a website, a brand, or a business plan.

We change that. We build complete businesses for underserved entrepreneurs in 2 to 5 days, at zero cost to them. And today, I want to show you exactly how this works, what your role is, and how we ALL benefit from building this together.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 2: THE PROBLEM - just add notes
# ═══════════════════════════════════════════════════════════════
notes(slides[1], """Here's the reality. Starting a business costs between 11 thousand and 54 thousand dollars before you even make a single sale.

For most people in our community, that's an impossible barrier. 72 percent of aspiring entrepreneurs never launch. Not because they don't have good ideas -- because they don't have the capital.

We eliminate every single one of these costs. We build the entire digital business infrastructure in 2 to 5 days, at zero cost to the entrepreneur.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 3: PROOF - ARLAN - just add notes
# ═══════════════════════════════════════════════════════════════
notes(slides[2], """Our first proof of concept. Arlan LLC -- premium exterior lighting and home services out of Greensboro.

We built their entire business from scratch in ONE day. Full market research. Professional brand. 8-page premium website. E-commerce store with 8 products. Live Stripe payments. AI-powered customer assistant.

Total market value: 41 thousand 500 dollars. One day. Zero cost to the entrepreneur.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 4: PROOF - HOOD HYMNS - just add notes
# ═══════════════════════════════════════════════════════════════
notes(slides[3], """Proof number two. Hood Hymns Publishing -- a full literary publishing platform for author DJ VanHook.

A Next.js application with trilingual support in English, Spanish, and Mandarin. AI-powered narration using Gemini text-to-speech. Cinematic book trailers. Merch store with live Stripe payments. Admin dashboard.

Total market value: 44 thousand dollars. Built in 2 days.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 5: PROOF - LAKE JEANETTE - just add notes
# ═══════════════════════════════════════════════════════════════
notes(slides[4], """Proof number three. Lake Jeanette Dentistry -- Brenes Precision Dentistry.

Professional dental practice website with service pages, team profiles, appointment system, mobile-responsive design, and local SEO optimization.

Total market value: 29 thousand dollars. One day.

Three businesses. Four days total. Over 114 thousand dollars in combined market value. Zero cost to the entrepreneurs. That's what Think! Ventures delivers.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 6: THE MODEL - add coaching, add notes
# ═══════════════════════════════════════════════════════════════
slide = slides[5]
# Add coaching card to existing model slide
card(slide, 0.8, 6.0, 11.7, 0.7, GOLD)
txt(slide, 1.1, 6.05, 11, 0.5,
    "NEW: 12-Month Coaching Program -- every entrepreneur receives mentoring to learn their systems and grow their business.",
    13, DARK_TEXT, True)

notes(slide, """Here's how the money flows.

On the left, Think! Ventures Foundation -- the nonprofit. Applies for grants, recruits entrepreneurs, runs workshops, measures impact.

On the right, Think! Design and Planning, LLC -- the for-profit company. Builds the technology.

The Foundation contracts with the LLC at fair market rates. 100 percent legal, 100 percent standard. Everything is disclosed on Form 990 and approved by the board.

And we've added something important -- a 12-month coaching program. Every entrepreneur we launch gets a full year of mentoring. Monthly check-ins. Help learning their website, their store, their financials. We don't just build the car and hand you the keys -- we teach you how to drive it.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 7: THE BOARD - REBUILD with new roles
# ═══════════════════════════════════════════════════════════════
slide = slides[6]
clear_slide(slide)

txt(slide, 0.8, 0.4, 5, 0.4, "YOUR BOARD", 14, GOLD, True)
txt(slide, 0.8, 0.9, 10, 0.6, "The Leadership Team", 36, TEAL, True)

members = [
    ("Chris Harrison", "Executive Director", "Founder, NC A&T faculty.\nStrategy, technology,\nprogram delivery."),
    ("Gregory Fulton", "Associate Director", "Operations and outreach.\nCommunity partnerships\nand recruitment."),
    ("Damon Howell", "Associate Director", "Operations and programs.\nWorkshop coordination\nand logistics."),
    ("Charlie Hopper", "Director (At-Large)", "Independent voice.\nCommunity perspective\nand accountability."),
    ("Cadasia Levy", "Director (At-Large)", "City of Greensboro.\nCoaching program lead.\nGovernment partnerships."),
]

for i, (name, role, desc) in enumerate(members):
    x = 0.3 + i * 2.55
    card(slide, x, 2.0, 2.35, 4.2)
    
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.65), Inches(2.2), Inches(1), Inches(1))
    circle.fill.solid()
    circle.fill.fore_color.rgb = CARD_BG
    circle.line.color.rgb = GOLD
    circle.line.width = Pt(2)
    txt(slide, x + 0.65, 2.4, 1, 0.5, "PHOTO", 10, MUTED_TEXT, False, PP_ALIGN.CENTER)
    
    txt(slide, x + 0.1, 3.4, 2.15, 0.4, name, 14, TEAL, True, PP_ALIGN.CENTER)
    txt(slide, x + 0.1, 3.8, 2.15, 0.3, role, 11, GOLD, True, PP_ALIGN.CENTER)
    txt(slide, x + 0.1, 4.2, 2.15, 1.5, desc, 10, MUTED_TEXT, False, PP_ALIGN.CENTER)

notes(slide, """This is our leadership team.

Chris Harrison -- Executive Director. Founder, NC A&T faculty. Handles strategy, technology, and program delivery.

Gregory Fulton -- Associate Director. Operations and outreach. Community partnerships and recruitment. 

Damon Howell -- Associate Director. Operations and programs. Workshop coordination and logistics.

Charlie Hopper -- Director at Large. Independent community voice. Brings accountability and fresh perspective.

Cadasia Levy -- Director at Large. Works for the City of Greensboro. She's leading our new coaching program and brings government partnership connections.

Five board members. Three in operational roles. Two independent voices. A Treasurer and Secretary will be appointed -- we're looking for a CPA to ensure our financials are handled by a professional.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 8: WHAT BOARD MEMBERS DO - just add notes
# ═══════════════════════════════════════════════════════════════
notes(slides[7], """There are TWO hats you can wear.

Board Service -- 4 meetings a year, about an hour each. You vote on budgets, review financials, approve contracts. That's the governance side.

Operational Work -- this is where you get paid. Community outreach, workshop facilitation, coaching, financial operations. Compensated as a contractor at fair market rates.

The bylaws specifically allow this. Each person recuses from the vote on their own compensation. Everything is transparent and on the books.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 9: COMPENSATION - REBUILD role-based
# ═══════════════════════════════════════════════════════════════
slide = slides[8]
clear_slide(slide)

txt(slide, 0.8, 0.4, 5, 0.4, "COMPENSATION", 14, GOLD, True)
txt(slide, 0.8, 0.9, 10, 0.6, "How Everyone Eats", 38, TEAL, True)

# Left card - Board service (everyone)
card(slide, 0.8, 2.0, 5.5, 4.2)
txt(slide, 1.1, 2.2, 5, 0.4, "All Board Members", 19, EMERALD, True)
bullets(slide, 1.1, 2.8, 5, 3.2, [
    "Meeting stipends: $100 - $250 per meeting",
    "4 meetings per year",
    "Approved in annual budget",
    "",
    "Every seat at this table has value.",
    "Every voice gets compensated for showing up."], 14, DARK_TEXT)

# Right card - Operational roles
card(slide, 7.0, 2.0, 5.5, 4.2)
txt(slide, 7.3, 2.2, 5, 0.4, "Operational Roles (As Needed)", 19, GOLD, True)
bullets(slide, 7.3, 2.8, 5, 3.2, [
    "Executive Director: part-time to full-time salary",
    "Associate Directors: $20 - $35/hr",
    "Technology services: LLC contract at market rates",
    "Coaching & mentoring: $20 - $35/hr",
    "Workshop facilitation: $20 - $35/hr",
    "Any member can take on operational work."], 14, DARK_TEXT)

notes(slide, """Here's how compensation works. Everyone at this table eats.

On the left -- Board Service. Every board member receives meeting stipends, 100 to 250 dollars per meeting, 4 meetings a year. That's for showing up, voting, reviewing financials. Every seat at this table has value.

On the right -- Operational Roles. These are separate from board duties. The Executive Director draws a salary that grows with grants. Associate Directors earn 20 to 35 dollars an hour for outreach, coordination, and logistics. My LLC gets paid for technology builds at fair market rates. Coaching and workshop facilitation -- also 20 to 35 dollars an hour.

Any board member can take on operational work. The board votes on each person's compensation with that person recused. Everything is transparent.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 10: TIMELINE - REBUILD role-based
# ═══════════════════════════════════════════════════════════════
slide = slides[9]
clear_slide(slide)

txt(slide, 0.8, 0.4, 5, 0.4, "THE TIMELINE", 14, GOLD, True)
txt(slide, 0.8, 0.9, 10, 0.6, "From Startup to Sustainability", 38, TEAL, True)

table(slide, 0.8, 2.0, 11.5, 4.2, [
    ["", "Phase 1: Pre-Grant", "Phase 2: First Grants", "Phase 3: Growth"],
    ["Foundation Revenue", "$0", "$58K - $165K", "$285K - $750K"],
    ["Executive Director", "LLC income only", "$30K - $60K/yr", "$80K - $150K/yr"],
    ["Associate Directors", "Volunteer (log hours)", "$15K - $35K/yr each", "$40K - $55K/yr each"],
    ["Board Stipends", "Waived", "$400 - $1,000/yr each", "$1,000 - $2,500/yr each"],
    ["Coaching Program", "Pilot (3 entrepreneurs)", "10 - 20 entrepreneurs", "50+ entrepreneurs"],
    ["Businesses Launched", "3 (proven)", "10 - 20", "50+"],
    ["Timeline", "Now", "Months 3 - 12", "Year 2 - 3"],
], col_widths=[2.8, 2.8, 3.0, 3.0])

notes(slide, """Three phases.

Phase 1 is right now. Pre-grant. The Executive Director lives off LLC income. Associate Directors volunteer and log their hours. Board stipends are waived until we're funded. We pilot the coaching program with our first 3 businesses.

Phase 2 -- first grants hit. Foundation revenue of 58 to 165 thousand. Executive Director draws 30 to 60 thousand. Associate Directors earn 15 to 35 thousand each. Board stipends kick in. Coaching program serves 10 to 20 entrepreneurs.

Phase 3 -- growth. 285 to 750 thousand in revenue. Everyone earning professional-level income. Coaching program at full scale with 50 or more entrepreneurs per year.

The key point -- log your hours now. When the grants arrive, those logged hours establish your track record and justify your compensation going forward.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 11: FUNDING SOURCES - just add notes
# ═══════════════════════════════════════════════════════════════
notes(slides[10], """Revenue comes from seven sources.

Federal grants -- EDA, SBA, MBDA. These are the big ones, 25 to 300 thousand depending on the year.

State grants -- NC IDEA, NC Commerce Department. 10 to 100 thousand.

Foundation grants -- Kauffman, Google for Startups. 5 to 75 thousand.

HBCU-specific grants -- HBCU Founders Initiative, Black Ambition. 10 to 75 thousand. Being connected to NC A&T gives us a strong position here.

Corporate sponsorships ramp up in Years 2 and 3.

Donations and workshop fees provide steady baseline revenue.

And the cooperative profit share -- every partner business shares 10 percent of profits with the Foundation. This grows as our ecosystem grows.

By Year 3, total budget of 285 to 750 thousand dollars.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 12: THE ASK - just add notes
# ═══════════════════════════════════════════════════════════════
notes(slides[11], """Here's the thing. I'm not asking you to do a hundred things.

If you believe in what you just saw -- the model, the proof, the vision -- then all I need is your signature on the bylaws.

I handle everything else. The IRS filing. The grant applications. The technology builds. The coaching program setup. The bank account. All of it.

You showed up. You believed. That's what matters. Sign the bylaws and let's get to work.""")

# ═══════════════════════════════════════════════════════════════
# SLIDE 13: VISION - just add notes
# ═══════════════════════════════════════════════════════════════
notes(slides[12], """Our vision is a world where lack of capital, technical knowledge, or professional connections never prevents a great idea from becoming a thriving business.

From dream to launch in days. We proved it with Arlan. We proved it with Hood Hymns. We proved it with Lake Jeanette Dentistry. We'll prove it 50 more times, then 500 more times.

And now with the coaching program, we're not just launching businesses -- we're growing entrepreneurs.

You're not just joining a board. You're joining a movement. Let's build something extraordinary together.

Now -- let's sign these documents and get to work.""")

# ═══════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════
prs.save(OUTPUT)
print(f"Presentation saved to: {OUTPUT}")
print(f"Total slides: {len(prs.slides)}")
print()
print("UPDATED:")
print("  - Slide 6: Added coaching program bar to the model")
print("  - Slide 7: Board rebuilt -- Associate Directors, Cadasia Levy added")
print("  - Slide 9: Compensation redesigned -- role-based, no names")
print("  - Slide 10: Timeline redesigned -- role-based + coaching row")
print("  - All slides: Speaker notes added")
print()
print("BOARD STRUCTURE:")
print("  Chris Harrison -- Executive Director")
print("  Gregory Fulton -- Associate Director")
print("  Damon Howell -- Associate Director")
print("  Charlie Hopper -- Director (At-Large)")
print("  Cadasia Levy -- Director (At-Large), City of Greensboro")
print("  Treasurer/Secretary -- TBD (CPA to be recruited)")
